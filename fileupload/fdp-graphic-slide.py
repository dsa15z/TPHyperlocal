#!/usr/bin/env python3
"""Generate a graphic, visual FDP sales slide — light theme, green accents, bold visuals."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: THE HERO — "What FDP Does" with visual pipeline
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
LIGHT_BG = RGBColor(0xEE, 0xF7, 0xED)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
DARK_GREEN = RGBColor(0x1B, 0x8A, 0x4A)
DEEP_GREEN = RGBColor(0x0D, 0x47, 0x2B)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x2D, 0x3A, 0x4A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
PURPLE = RGBColor(0x7D, 0x38, 0x8E)
CARD_SHADOW = RGBColor(0xE0, 0xE0, 0xE0)

# Stage colors
CLR_SOURCE = RGBColor(0xEF, 0x44, 0x44)
CLR_TRANSFORM = RGBColor(0xF5, 0x9E, 0x0B)
CLR_QUALITY = RGBColor(0x10, 0xB9, 0x81)
CLR_CDM = RGBColor(0x8B, 0x5C, 0xF6)
CLR_ENRICH = RGBColor(0x3B, 0x82, 0xF6)
CLR_RESOLVE = RGBColor(0xEC, 0x48, 0x99)
CLR_GOLD = RGBColor(0x06, 0xB6, 0xD4)

# Background
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = OFF_WHITE


def make_shadow():
    return {"type": "outer", "color": "000000", "blur": 10, "offset": 3, "angle": 135, "opacity": 0.08}


def add_shape(slide, shape_type, x, y, w, h, fill_color, line_color=None, shadow=False):
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if shadow:
        sp = s._element
        spPr = sp.find(qn('a:spPr')) or sp.find('.//' + qn('a:spPr'))
        if spPr is None:
            spPr = etree.SubElement(sp, qn('a:spPr'))
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', '127000')
        outerShdw.set('dist', '38100')
        outerShdw.set('dir', '8100000')
        outerShdw.set('algn', 'tl')
        srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgb.set('val', '000000')
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '12000')
    return s


def add_text(slide, text, x, y, w, h, size=12, color=BLACK, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font="Inter", valign=MSO_ANCHOR.TOP,
             spacing=None, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    if spacing:
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:spc'), str(int(spacing)))
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_rich_text(slide, parts, x, y, w, h, align=PP_ALIGN.LEFT):
    """parts = [(text, size, color, bold, font), ...]"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, color, bold, font in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return txBox


# ============================================================
# TOP BAR — Green accent
# ============================================================
add_shape(slide1, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, GREEN)

# ============================================================
# HEADER: FDP branding + headline
# ============================================================
add_text(slide1, "FDP", 0.7, 0.3, 1.5, 0.4, size=14, color=GRAY, bold=True, spacing=200)

add_rich_text(slide1, [
    ("THE COMPLETE ", 26, BLACK, True, "Arial"),
    ("DATA INTELLIGENCE ", 26, GREEN, True, "Arial"),
    ("PLATFORM", 26, BLACK, True, "Arial"),
], 0.7, 0.65, 11, 0.5)

add_text(slide1, "Everything Databricks makes you build, FDP ships out of the box.",
         0.7, 1.2, 7, 0.35, size=14, color=GRAY, italic=True)

# ============================================================
# LEFT SIDE: VISUAL PIPELINE — The Data Journey
# ============================================================

# Pipeline container card (white, rounded, with shadow)
pipeline_card = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
                          0.5, 1.8, 7.5, 4.8, WHITE, shadow=True)

# Pipeline title inside card
add_text(slide1, "THE DATA JOURNEY", 0.9, 1.95, 3, 0.3,
         size=10, color=GREEN, bold=True, spacing=250)
add_text(slide1, "From 5 raw sources to enterprise-grade gold tables",
         0.9, 2.2, 5, 0.25, size=9, color=GRAY)

# --- Pipeline Nodes (visual flow) ---
# Using large colored circles with labels, connected by lines

stages = [
    # (label, short_desc, color, stat)
    ("SOURCE", "5 data feeds", CLR_SOURCE, "2.9M"),
    ("TRANSFORM", "Normalize", CLR_TRANSFORM, ""),
    ("QUALITY", "Validate", CLR_QUALITY, "96%"),
    ("CDM", "Unify", CLR_CDM, "2.5M"),
    ("ENRICH", "Augment", CLR_ENRICH, ""),
    ("RESOLVE", "Deduplicate", CLR_RESOLVE, "1.4M"),
    ("GOLD", "Enterprise", CLR_GOLD, "99%"),
]

# Draw pipeline flow as a horizontal chain
circle_r = 0.42
start_x = 1.0
flow_y = 2.9
spacing_x = 0.95

for i, (label, desc, color, stat) in enumerate(stages):
    cx = start_x + i * spacing_x
    cy = flow_y

    # Connection line to next node
    if i < len(stages) - 1:
        line_x = cx + circle_r
        next_cx = start_x + (i + 1) * spacing_x
        line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(line_x), Inches(cy + circle_r / 2 - 0.015),
            Inches(next_cx - line_x), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE0, 0xE5, 0xEB)
        line.line.fill.background()

    # Circle node
    node = slide1.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx), Inches(cy), Inches(circle_r), Inches(circle_r))
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.fill.background()

    # Stat inside circle (if any)
    if stat:
        add_text(slide1, stat, cx - 0.05, cy + 0.08, circle_r + 0.1, circle_r - 0.1,
                 size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Label below
    add_text(slide1, label, cx - 0.18, cy + circle_r + 0.08, circle_r + 0.36, 0.2,
             size=7.5, color=DARK, bold=True, align=PP_ALIGN.CENTER, spacing=80)

    # Description below label
    add_text(slide1, desc, cx - 0.18, cy + circle_r + 0.27, circle_r + 0.36, 0.18,
             size=7, color=GRAY, align=PP_ALIGN.CENTER)

# --- Source labels (what feeds in) ---
sources = ["Salesforce CRM", "Dun & Bradstreet", "LinkedIn", "Web Crawlers", "SEC Filings"]
for i, src in enumerate(sources):
    y_pos = 4.05 + i * 0.2
    # Small colored dot
    dot = slide1.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(1.0), Inches(y_pos + 0.03), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CLR_SOURCE
    dot.line.fill.background()
    add_text(slide1, src, 1.15, y_pos, 1.5, 0.2, size=7.5, color=DARK)

# --- Enrichment labels ---
enrichments = ["Industry Classification", "Geocoding & Geo", "Revenue Estimation"]
for i, enr in enumerate(enrichments):
    y_pos = 4.05 + i * 0.2
    dot = slide1.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(3.2), Inches(y_pos + 0.03), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CLR_ENRICH
    dot.line.fill.background()
    add_text(slide1, enr, 3.35, y_pos, 1.5, 0.2, size=7.5, color=DARK)

# --- Output: Gold table result ---
gold_card = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
                      5.4, 4.05, 2.3, 0.95, LIGHT_BG)
add_text(slide1, "ENTERPRISE GOLD", 5.55, 4.12, 2.1, 0.22,
         size=8, color=DARK_GREEN, bold=True, spacing=100)
add_text(slide1, "1.4M unified records\n96% avg quality\n51% compression\n99% output accuracy",
         5.55, 4.35, 2.1, 0.6, size=8, color=DARK)


# ============================================================
# RIGHT SIDE: COMPARISON CALLOUT CARDS
# ============================================================

# "FDP vs Databricks" header
add_rich_text(slide1, [
    ("FDP ", 18, GREEN, True, "Arial"),
    ("vs ", 18, GRAY, False, "Arial"),
    ("DATABRICKS", 18, RGBColor(0xFF, 0x3E, 0x1D), True, "Arial"),
], 8.6, 1.85, 4.5, 0.4)

# What FDP includes that Databricks doesn't — green cards
fdp_wins = [
    ("Visual Data Journey", "See your entire pipeline\nend-to-end, not just code"),
    ("CRM Sync Built In", "Bidirectional Salesforce +\nHubSpot, no Fivetran needed"),
    ("Entity Resolution", "Match & merge across sources,\n2.5M down to 1.4M records"),
    ("AI Agent Tools", "MCP Registry with 27.6K\ninvocations — AI-native data"),
    ("Data Products", "Self-service marketplace with\nmedallion tiering & versioning"),
    ("Data Quality", "Great Expectations baked in,\nnot bolted on"),
]

card_w = 2.05
card_h = 0.82
card_gap = 0.1
cols = 2
start_cx = 8.5
start_cy = 2.35

for i, (title, desc) in enumerate(fdp_wins):
    col = i % cols
    row = i // cols
    cx = start_cx + col * (card_w + card_gap)
    cy = start_cy + row * (card_h + card_gap)

    # Green card
    card = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
                     cx, cy, card_w, card_h, GREEN, shadow=True)

    # Title
    add_text(slide1, title, cx + 0.12, cy + 0.08, card_w - 0.24, 0.24,
             size=10, color=WHITE, bold=True)

    # Description
    add_text(slide1, desc, cx + 0.12, cy + 0.34, card_w - 0.24, 0.44,
             size=8, color=WHITE)


# Databricks "requires" callout
db_y = start_cy + 3 * (card_h + card_gap) + 0.12
add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE,
          8.5, db_y, 4.2, 0.75, RGBColor(0xFF, 0xF3, 0xF0))

add_rich_text(slide1, [
    ("With Databricks ", 10, DARK, False, "Arial"),
    ("you'd need: ", 10, DARK, True, "Arial"),
    ("Fivetran + dbt + Monte Carlo + Atlan + custom entity resolution + a data engineering team.", 10, RGBColor(0xEF, 0x44, 0x44), False, "Arial"),
], 8.65, db_y + 0.1, 3.9, 0.6)


# ============================================================
# BOTTOM: Big stat callouts
# ============================================================
stats_y = 5.85
stat_data = [
    ("2.9M", "records in"),
    ("1.4M", "records out"),
    ("96%", "avg quality"),
    ("27.6K", "AI invocations"),
    ("94.6%", "pipeline success"),
]

stat_w = 2.2
for i, (val, label) in enumerate(stat_data):
    sx = 0.7 + i * stat_w
    # Big number
    add_text(slide1, val, sx, stats_y, stat_w, 0.45,
             size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, label, sx, stats_y + 0.45, stat_w, 0.25,
             size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# BOTTOM LINE
# ============================================================
add_rich_text(slide1, [
    ("// ", 10, GREEN, True, "Arial"),
    ("Databricks gives you a compute engine. ", 10, GRAY, False, "Arial"),
    ("FDP gives you the complete intelligence layer.", 10, BLACK, True, "Arial"),
], 0.7, 6.9, 12, 0.35)

# Footer
add_text(slide1, "futuridata.com", 11, 7.1, 2, 0.25,
         size=7, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/derekanderson/Documents/Personal (iCloud)/Claude/FDP-Graphic-Deck.pptx"
prs.save(output_path)
print(f"Done: {output_path}")
